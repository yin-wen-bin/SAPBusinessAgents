from __future__ import annotations

import asyncio
import json
import sys
from contextlib import contextmanager
from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
from fastapi.testclient import TestClient

from sap_business_agents_platform.app import create_app
from sap_business_agents_platform.config import Settings
from sap_business_agents_platform.codex_planner import CodexPlanner, _role_matching_thread_can_restart
from sap_business_agents_platform.database import RunStore
from sap_business_agents_platform.manifests import AgentRepository, is_agent_executable
from sap_business_agents_platform.role_matching import (
    ROLE_MATCHING_CATALOG_PAGE_CHARS,
    RoleMatchingService,
    _paginate_runtime_catalog,
)
from sap_business_agents_platform.role_matching_documents import load_chunks, preflight_scan, scan_and_extract


def _write_documents(root: Path) -> None:
    (root / "process.md").write_text("采购专员在SAP中创建采购订单。", encoding="utf-8")
    document = Document()
    document.add_heading("P2P", level=1)
    document.add_paragraph("应付会计核对供应商发票。")
    document.save(root / "process.docx")
    workbook = Workbook()
    workbook.active.title = "Operations"
    workbook.active.append(["岗位", "操作"])
    workbook.active.append(["销售", "查询销售订单交货状态"])
    workbook.save(root / "process.xlsx")
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text = "仓库主管核对收货"
    presentation.save(root / "process.pptx")
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with (root / "scanned.pdf").open("wb") as stream:
        writer.write(stream)
    (root / "legacy.doc").write_bytes(b"legacy")


def test_document_reader_tracks_locations_and_unsupported_files(tmp_path: Path) -> None:
    root = tmp_path / "documents"
    root.mkdir()
    _write_documents(root)
    preflight = preflight_scan(
        [str(root)], max_files=20, max_file_bytes=10_000_000, max_total_bytes=50_000_000
    )
    assert preflight["supported_file_count"] == 5
    assert preflight["ready"] is True
    result = scan_and_extract(
        [str(root)], cache_root=tmp_path / "cache", max_files=20,
        max_file_bytes=10_000_000, max_total_bytes=50_000_000,
    )
    parsed = {item["extension"]: item for item in result["documents"] if item["status"] == "parsed"}
    assert {".md", ".docx", ".xlsx", ".pptx"} <= set(parsed)
    assert load_chunks(parsed[".docx"])[0]["locator"]["kind"] == "docx_paragraph"
    assert load_chunks(parsed[".xlsx"])[0]["locator"]["sheet"] == "Operations"
    assert load_chunks(parsed[".pptx"])[0]["locator"]["slide"] == 1
    assert {item["code"] for item in result["issues"]} >= {
        "legacy_office_format_unsupported", "pdf_text_layer_unavailable"
    }
    assert result["extraction_complete"] is False


def test_incremental_document_scan_reuses_unchanged_parser_cache(tmp_path: Path) -> None:
    source = tmp_path / "process.txt"
    source.write_text("SAP收货操作", encoding="utf-8")
    first = scan_and_extract(
        [str(source)], cache_root=tmp_path / "cache", max_files=5,
        max_file_bytes=1000, max_total_bytes=1000,
    )
    reused = scan_and_extract(
        [str(source)], cache_root=tmp_path / "cache", max_files=5,
        max_file_bytes=1000, max_total_bytes=1000,
        reuse_by_hash={first["documents"][0]["sha256"]: first["documents"][0]},
    )
    assert reused["documents"][0]["reused"] is True


class _FakeRuntime:
    current_provider_id = "codex"

    def __init__(self) -> None:
        self.calls: list[dict] = []

    def supports(self, operation: str) -> bool:
        return operation in {"analyze_role_matching", "review_role_matching_feedback"}

    def snapshot(self, provider_id: str | None = None) -> dict:
        return {"provider_id": provider_id or "codex", "sdk_id": "codex-python-sdk", "version": "test", "configuration_digest": "test", "capabilities": ["role_matching"]}

    @contextmanager
    def pin(self, provider_id: str | None):
        yield

    async def cancel(self, thread_id: str | None = None) -> None:
        return None

    async def analyze_role_matching(self, **kwargs):
        self.calls.append(kwargs)
        document = kwargs["documents"][0]
        chunk = document["chunks"][0]
        ref = {"document_id": document["document_id"], "chunk_id": chunk["chunk_id"], "locator": chunk["locator"]}
        runtime_catalog = kwargs["agent_catalog"]["runtime_catalog"]
        evaluated_ids = [
            item["agent_id"] for page in runtime_catalog["pages"] for item in page["items"]
        ]
        return {
            "thread_id": "thread_role_test",
            "analysis": {
                "summary": {"zh": "识别到采购操作", "en": "Procurement operation identified"},
                "roles": [{"name": "采购专员", "evidence_refs": [ref]}],
                "processes": [{"name": "P2P", "evidence_refs": [ref]}],
                "operations": [{"operation_id": "op-1", "role": "采购专员", "department": "采购", "process": "P2P", "name": "查询采购订单", "description": "查询订单状态", "trigger": "日常", "inputs": [], "outputs": [], "sap_system_or_module": "MM", "frequency": "", "controls": [], "evidence_refs": [ref]}],
                "agent_matches": [{"operation_id": "op-1", "agent_id": "procure-to-pay-status", "coverage": "full", "confidence": "high", "reason": "流程一致", "uncovered_capabilities": [], "evidence_refs": [ref]}, {"operation_id": "op-1", "agent_id": "ar-collection", "coverage": "none", "confidence": "low", "reason": "业务范围不一致", "uncovered_capabilities": ["采购订单状态"], "evidence_refs": [ref]}, {"operation_id": "op-1", "agent_id": "invented-agent", "coverage": "full", "confidence": "high", "reason": "invalid", "evidence_refs": [ref]}],
                "workflow_suggestions": [], "agent_gaps": [], "document_issues": [],
                "non_sap_operation_count": 0,
                "catalog_evaluation": {
                    "catalog_digest": runtime_catalog["digest"],
                    "total_agent_count": runtime_catalog["total_agent_count"],
                    "evaluated_agent_count": len(evaluated_ids),
                    "evaluated_pair_count": len(evaluated_ids),
                    "evaluated_agent_ids": evaluated_ids,
                    "catalog_page_count": runtime_catalog["page_count"],
                    "agent_catalog_complete": True,
                    "matching_complete": True,
                    "failed_pages": [],
                },
            },
        }

    async def review_role_matching_feedback(self, **kwargs):
        return await self.analyze_role_matching(**kwargs)


class _Drafts:
    def create(self, *args, **kwargs):  # pragma: no cover - suggestions are empty here
        raise AssertionError("not expected")


async def _role_matching_session_scenario(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    selected = tmp_path / "selected"
    selected.mkdir()
    secret_text = "采购专员在SAP中查询采购订单 UNIQUE_BODY_74219"
    (selected / "p2p.txt").write_text(secret_text, encoding="utf-8")
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    store = RunStore(settings.database_path)
    runtime = _FakeRuntime()
    service = RoleMatchingService(settings, store, AgentRepository(repository_root / "agents"), runtime, _Drafts())
    await service.start()
    try:
        session = await service.create(
            paths=[str(selected)], role_description=None, locale="zh", consent=True
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"}:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        revision = service.revision(session["session_id"], 1)
        assert revision["result"]["agent_matches"][0]["agent_id"] == "procure-to-pay-status"
        assert len(revision["result"]["agent_matches"]) == 1
        assert revision["result"]["rejected_candidates"][0]["agent_id"] == "ar-collection"
        assert revision["result"]["catalog_evaluation"]["agent_catalog_complete"] is True
        assert runtime.calls and all("paths" not in document for document in runtime.calls[0]["documents"])
        assert secret_text.encode("utf-8") not in settings.database_path.read_bytes()
        assert all("path" not in json.dumps(event["data"]).lower() for event in store.role_matching_events_after(session["session_id"]))
        original_digest = revision["result_digest"]
        await service.feedback(
            session["session_id"], base_revision=1, message="聚焦采购岗位",
            mode="incremental", added_paths=[], added_role_description=None,
            excluded_document_ids=[],
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"} and current["current_revision"] == 2:
                break
            await asyncio.sleep(0.02)
        assert service.revision(session["session_id"], 1)["result_digest"] == original_digest
        assert service.revision(session["session_id"], 2)["parent_revision"] == 1
        assert service.documents(session["session_id"])[0].get("reused") is None  # internal hint is not contractual
        await service.feedback(
            session["session_id"], base_revision=2,
            message="使用当前完整Agent目录重新匹配", mode="full", added_paths=[],
            added_role_description=None, excluded_document_ids=[],
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"} and current["current_revision"] == 3:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        assert runtime.calls[-1]["reuse_business_understanding"] is True
        assert runtime.calls[-1]["previous_result"]["operations"]
    finally:
        await service.stop()


def test_role_matching_session_is_immutable_and_does_not_persist_body(tmp_path: Path) -> None:
    asyncio.run(_role_matching_session_scenario(tmp_path))


async def _description_source_scenario(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    description = "仓库主管在SAP中处理收货和盘点 UNIQUE_ROLE_TEXT_19384"
    supplemental = "月末核对库存差异并跟踪未清物料凭证 UNIQUE_ROLE_TEXT_62851"
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    store = RunStore(settings.database_path)
    runtime = _FakeRuntime()
    service = RoleMatchingService(
        settings, store, AgentRepository(repository_root / "agents"), runtime, _Drafts()
    )
    await service.start()
    try:
        session = await service.create(
            paths=[], role_description=description, locale="zh", consent=True
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"}:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        sources = service.documents(session["session_id"])
        assert len(sources) == 1
        assert sources[0]["source_type"] == "user_description"
        assert sources[0]["paths"] == []
        assert runtime.calls[0]["documents"][0]["source_type"] == "user_description"
        result = service.revision(session["session_id"], 1)["result"]
        assert result["completeness"]["document_scan_status"] == "not_requested"
        assert result["completeness"]["document_source_count"] == 0
        assert result["completeness"]["description_source_count"] == 1
        assert result["agent_matches"][0]["evidence_refs"][0]["source_type"] == "user_description"
        assert description.encode("utf-8") not in settings.database_path.read_bytes()
        assert description not in service.markdown(session["session_id"], 1)
        csv_report = service.csv(session["session_id"], 1, "agent_matches")
        assert "source_types" in csv_report
        assert "user_description" in csv_report
        assert description not in csv_report

        combined_file = tmp_path / "combined.txt"
        combined_file.write_text("采购岗位在SAP中查询采购订单", encoding="utf-8")
        combined = await service.create(
            paths=[str(combined_file)], role_description="采购岗位还负责跟踪收货",
            locale="zh", consent=True,
        )
        for _ in range(100):
            combined_current = service.get(combined["session_id"])
            if combined_current["status"] in {"completed", "failed"}:
                break
            await asyncio.sleep(0.02)
        assert combined_current["status"] == "completed", combined_current.get("error")
        assert {
            item["source_type"] for item in runtime.calls[-1]["documents"]
        } == {"document", "user_description"}

        await service.feedback(
            session["session_id"], base_revision=1, message="补充月结职责",
            mode="incremental", added_paths=[],
            added_role_description=supplemental, excluded_document_ids=[],
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"} and current["current_revision"] == 2:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        assert len(service.documents(session["session_id"])) == 2
        assert supplemental.encode("utf-8") not in settings.database_path.read_bytes()

        calls_before = len(runtime.calls)
        source_ids = [item["document_id"] for item in service.documents(session["session_id"])]
        await service.feedback(
            session["session_id"], base_revision=2, message="排除全部来源",
            mode="full", added_paths=[], added_role_description=None,
            excluded_document_ids=source_ids,
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] == "failed":
                break
            await asyncio.sleep(0.02)
        assert current["error"]["code"] == "role_matching_no_active_sources"
        assert len(runtime.calls) == calls_before
    finally:
        await service.stop()


def test_role_description_is_a_citable_private_source(tmp_path: Path) -> None:
    asyncio.run(_description_source_scenario(tmp_path))


def test_role_matching_assistant_is_not_executable_or_composable() -> None:
    root = Path(__file__).resolve().parents[1]
    manifest = AgentRepository(root / "agents").get("role-agent-matching")
    assert manifest["kind"] == "platform_assistant"
    assert manifest["assistant"]["composable"] is False
    assert is_agent_executable(manifest) is False


def test_archived_role_matching_thread_can_restart_but_other_runtime_errors_cannot() -> None:
    assert _role_matching_thread_can_restart(
        RuntimeError("JSON-RPC error -32600: session abc is archived")
    ) is True
    assert _role_matching_thread_can_restart(RuntimeError("authentication failed")) is False


class _RetryAfterFailureRuntime(_FakeRuntime):
    def __init__(self) -> None:
        super().__init__()
        self.feedback_attempts = 0

    async def review_role_matching_feedback(self, **kwargs):
        self.feedback_attempts += 1
        if self.feedback_attempts == 1:
            raise RuntimeError("temporary Runtime failure")
        return await self.analyze_role_matching(**kwargs)


async def _failed_role_matching_turn_can_retry_scenario(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    store = RunStore(settings.database_path)
    runtime = _RetryAfterFailureRuntime()
    service = RoleMatchingService(
        settings, store, AgentRepository(repository_root / "agents"), runtime, _Drafts(),
    )
    await service.start()
    try:
        session = await service.create(
            paths=[], role_description="查询销售订单到收款状态", locale="zh", consent=True,
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"}:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed"

        failed_attempt = await service.feedback(
            session["session_id"], base_revision=1, message="全量重新匹配",
            mode="full", added_paths=[], added_role_description=None,
            excluded_document_ids=[],
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] == "failed":
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "failed"
        assert current["active_job_id"] is None
        assert current["current_revision"] == 1

        # Older processes could leave the ID of an already failed scheduler
        # job on the session.  That terminal reference must not block retry.
        store.update_role_matching_session(
            session["session_id"], active_job_id=failed_attempt["active_job_id"]
        )

        await service.feedback(
            session["session_id"], base_revision=1, message="再次全量重新匹配",
            mode="full", added_paths=[], added_role_description=None,
            excluded_document_ids=[],
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"} and current["current_revision"] == 2:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        assert current["current_revision"] == 2
        assert service.revision(session["session_id"], 1)["result_digest"]
        assert service.revision(session["session_id"], 2)["parent_revision"] == 1
    finally:
        await service.stop()


def test_failed_role_matching_turn_can_retry_last_valid_revision(tmp_path: Path) -> None:
    asyncio.run(_failed_role_matching_turn_can_retry_scenario(tmp_path))


def test_codex_role_matching_uses_isolated_catalog_pages_and_compact_coverage_proof(
    monkeypatch: pytest.MonkeyPatch, tmp_path: Path,
) -> None:
    started_threads: list[object] = []
    page_attempts: dict[str, int] = {}

    def response(analysis: dict) -> SimpleNamespace:
        return SimpleNamespace(
            final_response=json.dumps(
                {
                    "analysis_json": json.dumps(analysis, ensure_ascii=False),
                    "summary_zh": "完成",
                    "summary_en": "Complete",
                },
                ensure_ascii=False,
            )
        )

    class FakeThread:
        def __init__(self, index: int) -> None:
            self.id = f"thread-{index}"

        async def run(self, prompt: str, **_kwargs):
            if prompt.startswith("Analyze"):
                return response(
                    {
                        "roles": [], "processes": [], "document_issues": [],
                        "operations": [{"operation_id": "op-1", "evidence_refs": []}],
                        "non_sap_operation_count": 0,
                    }
                )
            if prompt.startswith("Evaluate"):
                agent_id = "agent-a" if '"agent-a"' in prompt else "agent-b"
                page_attempts[agent_id] = page_attempts.get(agent_id, 0) + 1
                if agent_id == "agent-b" and page_attempts[agent_id] == 1:
                    raise RuntimeError("transient catalog-page interruption")
                coverage = "partial" if agent_id == "agent-a" else "none"
                candidate = {
                    "operation_id": "op-1", "agent_id": agent_id,
                    "coverage": coverage, "confidence": "medium", "reason": "test",
                    "uncovered_capabilities": [], "evidence_refs": [],
                }
                return response(
                    {
                        "roles": [], "processes": [], "operations": [],
                        "agent_matches": [candidate] if coverage != "none" else [],
                        "rejected_candidates": [candidate] if coverage == "none" else [],
                        "workflow_suggestions": [], "agent_gaps": [], "document_issues": [],
                        "catalog_evaluation": {
                            "catalog_digest": "sha256:test", "total_agent_count": 2,
                            "evaluated_agent_count": 1, "evaluated_pair_count": 1,
                            "evaluated_agent_ids": [agent_id], "catalog_page_count": 2,
                            "agent_catalog_complete": True, "matching_complete": True,
                            "failed_pages": [],
                        },
                    }
                )
            assert prompt.startswith("Finalize")
            return response({"workflow_suggestions": [], "agent_gaps": []})

    class FakeCodex:
        async def __aenter__(self):
            return self

        async def __aexit__(self, *_args):
            return None

        async def thread_start(self, **_kwargs):
            thread = FakeThread(len(started_threads) + 1)
            started_threads.append(thread)
            return thread

    monkeypatch.setitem(
        sys.modules,
        "openai_codex",
        SimpleNamespace(
            AsyncCodex=FakeCodex,
            ApprovalMode=SimpleNamespace(deny_all="deny_all"),
            Sandbox=SimpleNamespace(read_only="read_only"),
        ),
    )
    catalog = {
        "runtime_catalog": {
            "digest": "sha256:test", "total_agent_count": 2, "page_count": 2,
            "pages": [
                {"catalog_digest": "sha256:test", "page_index": 1, "page_count": 2,
                 "total_agent_count": 2, "items": [{"agent_id": "agent-a"}]},
                {"catalog_digest": "sha256:test", "page_index": 2, "page_count": 2,
                 "total_agent_count": 2, "items": [{"agent_id": "agent-b"}]},
            ],
        }
    }
    result = asyncio.run(
        CodexPlanner(tmp_path).analyze_role_matching(
            documents=[], agent_catalog=catalog, previous_result=None, user_context="",
            rematch_mode="full", locale="zh", thread_id=None,
        )
    )["analysis"]
    assert len(started_threads) == 4  # primary, page A, failed page B, retried page B
    assert page_attempts == {"agent-a": 1, "agent-b": 2}
    assert result["catalog_evaluation"]["agent_catalog_complete"] is True
    assert result["catalog_evaluation"]["evaluated_pair_count"] == 2
    assert result["agent_matches"][0]["agent_id"] == "agent-a"
    assert result["rejected_candidates"][0]["agent_id"] == "agent-b"


def test_runtime_catalog_is_complete_paged_json_without_duplicate_compiler_catalog(
    tmp_path: Path,
) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    service = RoleMatchingService(
        settings, RunStore(settings.database_path),
        AgentRepository(repository_root / "agents"), _FakeRuntime(), _Drafts(),
    )
    catalog = service._catalog()
    runtime_catalog = catalog["runtime_catalog"]
    assert "executable_catalog" not in runtime_catalog
    ids = []
    for page in runtime_catalog["pages"]:
        encoded = json.dumps(page, ensure_ascii=False, separators=(",", ":"))
        assert len(encoded) <= ROLE_MATCHING_CATALOG_PAGE_CHARS
        assert "…[TRUNCATED]" not in encoded
        assert json.loads(encoded)["catalog_digest"] == catalog["digest"]
        ids.extend(item["agent_id"] for item in page["items"])
    assert len(ids) == runtime_catalog["total_agent_count"]
    assert len(ids) == len(set(ids))
    assert "ar-collection" in ids
    assert "order-to-cash-status" in ids
    o2c = next(
        item for page in runtime_catalog["pages"] for item in page["items"]
        if item["agent_id"] == "order-to-cash-status"
    )
    assert "pgi_status" in o2c["capability_signals"]

    # A synthetic catalog larger than the former 120k character cutoff must
    # still be split only at complete Agent object boundaries, including when
    # the relevant O2C Agent is deliberately placed last.
    synthetic = [
        {"agent_id": f"agent-{index:03d}", "summary": "x" * 5000}
        for index in range(25)
    ]
    synthetic.append({"agent_id": "order-to-cash-status", "summary": "x" * 5000})
    pages = _paginate_runtime_catalog(
        synthetic, digest="sha256:synthetic", max_chars=45_000
    )
    serialized = [
        json.dumps(page, ensure_ascii=False, separators=(",", ":")) for page in pages
    ]
    assert sum(len(item) for item in serialized) > 120_000
    assert all(len(item) <= 45_000 and "…[TRUNCATED]" not in item for item in serialized)
    assert pages[-1]["items"][-1]["agent_id"] == "order-to-cash-status"


class _O2CRuntime(_FakeRuntime):
    async def analyze_role_matching(self, **kwargs):
        self.calls.append(kwargs)
        runtime_catalog = kwargs["agent_catalog"]["runtime_catalog"]
        catalog_ids = [
            item["agent_id"] for page in runtime_catalog["pages"] for item in page["items"]
        ]
        assert "order-to-cash-status" in catalog_ids
        document = kwargs["documents"][0]
        chunk = document["chunks"][0]
        ref = {
            "document_id": document["document_id"], "chunk_id": chunk["chunk_id"],
            "locator": chunk["locator"],
        }
        operation = {
            "operation_id": "op-o2c", "role": "销售订单到收款状态查询人员",
            "department": "销售", "process": "O2C", "name": "查询订单到收款状态",
            "description": "查询销售订单、交货、PGI、开票和FI清账状态",
            "trigger": "输入销售订单号", "inputs": ["销售订单号"],
            "outputs": ["O2C状态"], "sap_system_or_module": "SD/FI-AR",
            "frequency": "", "controls": ["只读"], "evidence_refs": [ref],
        }
        return {
            "thread_id": "thread_o2c_test",
            "analysis": {
                "summary": {"zh": "已匹配O2C", "en": "O2C matched"},
                "roles": [{"name": operation["role"], "evidence_refs": [ref]}],
                "processes": [{"name": "O2C", "evidence_refs": [ref]}],
                "operations": [operation],
                "agent_matches": [
                    {"operation_id": "op-o2c", "agent_id": "ar-collection", "coverage": "partial", "confidence": "medium", "reason": "仅覆盖客户维度应收催收", "uncovered_capabilities": ["销售订单到交货和开票追踪"], "evidence_refs": [ref]},
                    {"operation_id": "op-o2c", "agent_id": "order-to-cash-status", "coverage": "partial", "confidence": "low", "reason": "按销售订单追踪完整O2C状态", "uncovered_capabilities": ["明确PGI状态判定"], "evidence_refs": [ref]},
                ],
                "rejected_candidates": [], "workflow_suggestions": [], "agent_gaps": [{
                    "gap_id": "false-o2c-gap", "operation_ids": ["op-o2c"],
                    "required_capability": "错误的O2C缺口", "required_inputs": ["sales_order"],
                    "required_outputs": ["business_status"], "safety_boundary": "只读",
                    "business_impact": "无", "partial_agent_ids": [],
                    "reason": "完整匹配存在时不应保留", "evidence_refs": [ref],
                }],
                "document_issues": [], "non_sap_operation_count": 0,
                "catalog_evaluation": {
                    "catalog_digest": runtime_catalog["digest"],
                    "total_agent_count": len(catalog_ids),
                    "evaluated_agent_count": len(catalog_ids),
                    "evaluated_pair_count": len(catalog_ids),
                    "evaluated_agent_ids": catalog_ids,
                    "catalog_page_count": runtime_catalog["page_count"],
                    "agent_catalog_complete": True, "matching_complete": True,
                    "consolidation_complete": True,
                    "failed_pages": [],
                },
            },
        }


async def _o2c_complete_catalog_scenario(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    service = RoleMatchingService(
        settings, RunStore(settings.database_path), AgentRepository(repository_root / "agents"),
        _O2CRuntime(), _Drafts(),
    )
    await service.start()
    try:
        session = await service.create(
            paths=[],
            role_description="负责查询销售订单、交货、PGI、开票和应收清账状态。",
            locale="zh", consent=True,
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"}:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        result = service.revision(session["session_id"], 1)["result"]
        assert [item["agent_id"] for item in result["agent_matches"]] == [
            "order-to-cash-status", "ar-collection"
        ]
        assert result["agent_matches"][0]["confidence"] == "high"
        assert result["agent_gaps"] == []
        assert result["completeness"]["agent_catalog_complete"] is True
        assert result["completeness"]["matching_complete"] is True
        assert result["completeness"]["workflow_validation_complete"] is True
    finally:
        await service.stop()


def test_o2c_role_matching_prefers_order_to_cash_status_with_complete_catalog(
    tmp_path: Path,
) -> None:
    asyncio.run(_o2c_complete_catalog_scenario(tmp_path))


class _IncompleteCatalogRuntime(_FakeRuntime):
    async def analyze_role_matching(self, **kwargs):
        result = await super().analyze_role_matching(**kwargs)
        analysis = result["analysis"]
        evaluation = analysis["catalog_evaluation"]
        evaluation["evaluated_agent_ids"] = evaluation["evaluated_agent_ids"][:-1]
        evaluation["evaluated_agent_count"] -= 1
        evaluation["evaluated_pair_count"] -= 1
        evaluation["agent_catalog_complete"] = False
        evaluation["matching_complete"] = False
        evaluation["failed_pages"] = [evaluation["catalog_page_count"]]
        ref = analysis["operations"][0]["evidence_refs"]
        analysis["agent_gaps"] = [{
            "gap_id": "unsafe-gap", "operation_ids": ["op-1"],
            "required_capability": "未经完整目录确认的缺口", "required_inputs": [],
            "required_outputs": [], "safety_boundary": "只读", "business_impact": "未知",
            "partial_agent_ids": [], "reason": "目录不完整", "evidence_refs": ref,
        }]
        return result


async def _incomplete_catalog_scenario(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    service = RoleMatchingService(
        settings, RunStore(settings.database_path), AgentRepository(repository_root / "agents"),
        _IncompleteCatalogRuntime(), _Drafts(),
    )
    await service.start()
    try:
        session = await service.create(
            paths=[], role_description="负责在SAP中查询采购订单", locale="zh", consent=True
        )
        for _ in range(100):
            current = service.get(session["session_id"])
            if current["status"] in {"completed", "failed"}:
                break
            await asyncio.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        result = service.revision(session["session_id"], 1)["result"]
        assert result["catalog_evaluation"]["agent_catalog_complete"] is False
        assert result["completeness"]["matching_complete"] is False
        assert result["agent_gaps"] == []
        assert result["workflow_suggestions"] == []
        assert result["workflow_validation_issues"] == [
            {"code": "role_matching_catalog_incomplete"}
        ]
    finally:
        await service.stop()


def test_incomplete_catalog_suppresses_gaps_and_workflow_suggestions(
    tmp_path: Path,
) -> None:
    asyncio.run(_incomplete_catalog_scenario(tmp_path))


def test_workflow_suggestions_use_only_executable_agents_and_compiler_v4(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts",
    )
    service = RoleMatchingService(
        settings,
        RunStore(settings.database_path),
        AgentRepository(repository_root / "agents"),
        _FakeRuntime(),
        _Drafts(),
    )
    catalog = service._catalog()
    valid = {
        "suggestion_id": "wf-p2p",
        "title": {"zh": "采购订单跟踪", "en": "Purchase order tracking"},
        "description": "Trace purchase-order business status.",
        "stages": [
            {
                "id": "trace_purchase_order",
                "capability": {"zh": "跟踪采购订单", "en": "Trace purchase order"},
                "agent_id": "procure-to-pay-status",
                "confidence": "high",
                "reason": "The agent provides the required P2P evidence chain.",
                "bindings": [],
                "requested_outputs": [
                    "po_results", "business_status", "source_complete",
                    "evidence_complete", "business_report",
                ],
            }
        ],
    }
    compiled, issues = service._compile_suggestions(
        "rolesession_12345678", [valid], catalog, "zh"
    )
    assert issues == []
    assert compiled[0]["compiler_version"] == 4
    assert compiled[0]["compiled_workflow"]["readOnly"] is True
    assert compiled[0]["compiled_workflow"]["nodes"][0]["agentId"] == "procure-to-pay-status"

    blocked = {
        **valid,
        "suggestion_id": "wf-blocked",
        "stages": [{**valid["stages"][0], "agent_id": "billing-output-monitor"}],
    }
    rejected, blocked_issues = service._compile_suggestions(
        "rolesession_12345678", [blocked], catalog, "zh"
    )
    assert rejected == []
    assert blocked_issues


def test_role_matching_api_preflight_and_session(tmp_path: Path) -> None:
    repository_root = Path(__file__).resolve().parents[1]
    source = tmp_path / "role.txt"
    source.write_text("SAP采购专员查询采购订单", encoding="utf-8")
    settings = replace(
        Settings(), repository_root=repository_root, data_root=tmp_path / "data",
        draft_root=tmp_path / "drafts", free_query_runtime="planner_legacy",
    )
    app = create_app(settings, planner=_FakeRuntime())
    with TestClient(app) as client:
        preflight = client.post("/api/role-matching/preflight", json={"paths": [str(source)]})
        assert preflight.status_code == 200
        assert preflight.json()["ready"] is True
        description_preflight = client.post(
            "/api/role-matching/preflight",
            json={"paths": [], "roleDescription": "仓库主管负责SAP库存盘点"},
        )
        assert description_preflight.status_code == 200
        assert description_preflight.json()["source_mode"] == "description"
        assert description_preflight.json()["supported_file_count"] == 0
        assert client.post(
            "/api/role-matching/sessions",
            json={"paths": [], "roleDescription": "   ", "locale": "zh", "consentToRuntime": True},
        ).status_code == 422
        oversized_description = "PRIVATE_DESCRIPTION_SHOULD_NOT_LEAK_" + "x" * 12_000
        oversized_response = client.post(
            "/api/role-matching/sessions",
            json={"paths": [], "roleDescription": oversized_description, "locale": "zh", "consentToRuntime": True},
        )
        assert oversized_response.status_code == 422
        assert "PRIVATE_DESCRIPTION_SHOULD_NOT_LEAK" not in oversized_response.text
        rejected = client.post(
            "/api/role-matching/sessions",
            json={"paths": [str(source)], "locale": "zh", "consentToRuntime": False},
        )
        assert rejected.status_code == 422
        response = client.post(
            "/api/role-matching/sessions",
            json={"paths": [str(source)], "locale": "zh", "consentToRuntime": True},
        )
        assert response.status_code == 202, response.text
        session_id = response.json()["session_id"]
        for _ in range(100):
            current = client.get(f"/api/role-matching/sessions/{session_id}").json()
            if current["status"] in {"completed", "failed"}:
                break
            import time
            time.sleep(0.02)
        assert current["status"] == "completed", current.get("error")
        assert client.get(f"/api/role-matching/sessions/{session_id}/revisions/1").status_code == 200
        assert client.get(f"/api/role-matching/sessions/{session_id}/revisions/1/report.md").status_code == 200
        json_report = client.get(
            f"/api/role-matching/sessions/{session_id}/revisions/1/report.json"
        )
        assert json_report.status_code == 200
        assert json_report.headers["content-disposition"].endswith(
            'filename="role-matching-report.json"'
        )
