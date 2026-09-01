from __future__ import annotations

import csv
import hashlib
import json
import os
from pathlib import Path
from typing import Any, Iterable


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".yaml",
    ".yml",
}
LEGACY_EXTENSIONS = {".doc", ".xls", ".ppt"}
IMAGE_EXTENSIONS = {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tif", ".tiff"}
PARSER_VERSION = "role-documents.v1"
USER_DESCRIPTION_EXTENSION = ".user-description"
USER_DESCRIPTION_PARSER_VERSION = "role-user-description.v1"


class RoleDocumentError(ValueError):
    def __init__(self, message: str, *, code: str, detail: Any = None) -> None:
        super().__init__(message)
        self.code = code
        self.detail = detail


def preflight_scan(
    roots: list[str], *, max_files: int, max_file_bytes: int, max_total_bytes: int
) -> dict[str, Any]:
    resolved = _resolve_roots(roots)
    candidates, issues = _discover(resolved)
    sizes: list[int] = []
    for path in candidates:
        try:
            size = path.stat().st_size
        except OSError:
            issues.append(_issue(path, "file_unavailable"))
            continue
        sizes.append(size)
        if size > max_file_bytes:
            issues.append(_issue(path, "file_too_large", size=size))
    total = sum(sizes)
    blockers = []
    if len(candidates) > max_files:
        blockers.append({"code": "role_matching_file_limit", "count": len(candidates), "limit": max_files})
    if total > max_total_bytes:
        blockers.append({"code": "role_matching_total_size_limit", "bytes": total, "limit": max_total_bytes})
    if any(item["code"] in {"file_too_large", "file_unavailable"} for item in issues):
        blockers.append({"code": "role_matching_document_preflight_failed"})
    return {
        "roots": [str(path) for path in resolved], "supported_file_count": len(candidates),
        "total_bytes": total, "issues": issues, "blockers": blockers,
        "ready": not blockers,
    }


def scan_and_extract(
    roots: list[str],
    *,
    cache_root: Path,
    max_files: int,
    max_file_bytes: int,
    max_total_bytes: int,
    reuse_by_hash: dict[str, dict[str, Any]] | None = None,
) -> dict[str, Any]:
    resolved_roots = _resolve_roots(roots)
    candidates, issues = _discover(resolved_roots)
    if len(candidates) > max_files:
        raise RoleDocumentError(
            f"The selected paths contain {len(candidates)} supported documents; the limit is {max_files}.",
            code="role_matching_file_limit",
            detail={"count": len(candidates), "limit": max_files},
        )
    try:
        total_bytes = sum(path.stat().st_size for path in candidates)
    except OSError as exc:
        raise RoleDocumentError(
            "A selected document became unavailable during scanning.",
            code="role_matching_file_unavailable",
            detail={"message": str(exc)},
        ) from exc
    if total_bytes > max_total_bytes:
        raise RoleDocumentError(
            "The selected documents exceed the total-size limit.",
            code="role_matching_total_size_limit",
            detail={"bytes": total_bytes, "limit": max_total_bytes},
        )

    cache_root.mkdir(parents=True, exist_ok=True)
    by_hash: dict[str, dict[str, Any]] = {}
    for path in candidates:
        before = path.stat()
        if before.st_size > max_file_bytes:
            issues.append(_issue(path, "file_too_large", size=before.st_size))
            continue
        digest = _sha256(path)
        after = path.stat()
        if (before.st_size, before.st_mtime_ns) != (after.st_size, after.st_mtime_ns):
            issues.append(_issue(path, "file_changed_during_read"))
            continue
        existing = by_hash.get(digest)
        if existing is not None:
            existing["paths"].append(str(path))
            continue
        document_id = f"doc_{digest[:20]}"
        cache_path = cache_root / f"{document_id}.json"
        reusable = (reuse_by_hash or {}).get(digest)
        reusable_cache = Path(str((reusable or {}).get("cache_path") or ""))
        if (
            reusable
            and reusable.get("parser_version") == PARSER_VERSION
            and reusable_cache.is_file()
        ):
            record = {
                **{key: value for key, value in reusable.items() if key in {
                    "document_id", "sha256", "name", "extension", "size_bytes", "modified_ns",
                    "parser_version", "status", "issue", "chunk_count", "cache_path",
                }},
                "name": path.name,
                "size_bytes": before.st_size,
                "modified_ns": before.st_mtime_ns,
                "paths": [str(path)],
                "reused": True,
            }
            by_hash[digest] = record
            continue
        extracted = _extract(path)
        record = {
            "document_id": document_id,
            "sha256": digest,
            "name": path.name,
            "extension": path.suffix.lower(),
            "size_bytes": before.st_size,
            "modified_ns": before.st_mtime_ns,
            "paths": [str(path)],
            "parser_version": PARSER_VERSION,
            "status": extracted["status"],
            "issue": extracted.get("issue"),
            "chunk_count": len(extracted.get("chunks") or []),
            "cache_path": str(cache_path),
            "reused": False,
        }
        cache_path.write_text(
            json.dumps(
                {
                    "document_id": document_id,
                    "parser_version": PARSER_VERSION,
                    "chunks": extracted.get("chunks") or [],
                },
                ensure_ascii=False,
                separators=(",", ":"),
            ),
            encoding="utf-8",
        )
        if record["status"] != "parsed":
            issues.append(
                {
                    "document_id": document_id,
                    "name": path.name,
                    "code": str(record.get("issue") or "document_unparseable"),
                }
            )
        by_hash[digest] = record
    return {
        "roots": [str(path) for path in resolved_roots],
        "documents": list(by_hash.values()),
        "issues": issues,
        "file_count": len(candidates),
        "unique_document_count": len(by_hash),
        "total_bytes": total_bytes,
        "scan_complete": not any(item["code"] == "file_changed_during_read" for item in issues),
        "extraction_complete": (
            all(item["status"] == "parsed" for item in by_hash.values())
            and not issues
        ),
    }


def load_chunks(document: dict[str, Any]) -> list[dict[str, Any]]:
    cache_path = Path(str(document["cache_path"]))
    payload = json.loads(cache_path.read_text(encoding="utf-8"))
    return list(payload.get("chunks") or [])


def create_user_description_document(
    text: str,
    *,
    cache_root: Path,
    turn: int,
    locale: str,
) -> dict[str, Any]:
    content = str(text).strip()
    if not content:
        raise RoleDocumentError(
            "Role description cannot be blank.",
            code="role_matching_description_blank",
        )
    content_digest = hashlib.sha256(content.encode("utf-8")).hexdigest()
    identity_digest = hashlib.sha256(
        f"{turn}\n{content}".encode("utf-8")
    ).hexdigest()
    document_id = f"role_text_{identity_digest[:20]}"
    cache_root.mkdir(parents=True, exist_ok=True)
    cache_path = cache_root / f"{document_id}.json"
    label = (
        f"用户提供的岗位描述 · 第{turn}轮"
        if locale == "zh"
        else f"User-provided role description · Turn {turn}"
    )
    chunks = [
        _chunk(
            content,
            {"kind": "user_description", "turn": turn, "label": label},
        )
    ]
    cache_path.write_text(
        json.dumps(
            {
                "document_id": document_id,
                "parser_version": USER_DESCRIPTION_PARSER_VERSION,
                "chunks": chunks,
            },
            ensure_ascii=False,
            separators=(",", ":"),
        ),
        encoding="utf-8",
    )
    return {
        "document_id": document_id,
        "sha256": content_digest,
        "name": label,
        "extension": USER_DESCRIPTION_EXTENSION,
        "size_bytes": len(content.encode("utf-8")),
        "modified_ns": "",
        "paths": [],
        "parser_version": USER_DESCRIPTION_PARSER_VERSION,
        "status": "parsed",
        "issue": None,
        "chunk_count": len(chunks),
        "cache_path": str(cache_path),
        "reused": False,
        "source_type": "user_description",
    }


def empty_document_scan() -> dict[str, Any]:
    return {
        "roots": [],
        "documents": [],
        "issues": [],
        "file_count": 0,
        "unique_document_count": 0,
        "total_bytes": 0,
        "scan_complete": True,
        "extraction_complete": True,
        "scan_status": "not_requested",
    }


def _resolve_roots(roots: list[str]) -> list[Path]:
    resolved: list[Path] = []
    seen: set[str] = set()
    for raw in roots:
        value = str(raw).strip()
        if not value:
            raise RoleDocumentError(
                "Document paths cannot be blank.", code="role_matching_path_blank"
            )
        path = Path(value)
        if not path.is_absolute():
            raise RoleDocumentError(
                f"Document path must be absolute: {value}",
                code="role_matching_path_not_absolute",
            )
        try:
            canonical = path.resolve(strict=True)
        except OSError as exc:
            raise RoleDocumentError(
                f"Document path is unavailable: {value}",
                code="role_matching_path_unavailable",
                detail={"message": str(exc)},
            ) from exc
        key = os.path.normcase(str(canonical))
        if key not in seen:
            seen.add(key)
            resolved.append(canonical)
    if not resolved:
        raise RoleDocumentError(
            "At least one document path is required.", code="role_matching_paths_required"
        )
    return resolved


def _discover(roots: Iterable[Path]) -> tuple[list[Path], list[dict[str, Any]]]:
    files: list[Path] = []
    issues: list[dict[str, Any]] = []
    seen: set[str] = set()
    for root in roots:
        if root.is_file():
            candidates = [root]
        else:
            candidates = []
            for directory, names, filenames in os.walk(root, followlinks=False):
                base = Path(directory)
                names[:] = [
                    name for name in names if not _is_reparse_or_link(base / name)
                ]
                candidates.extend(base / name for name in filenames)
        for path in candidates:
            if _is_reparse_or_link(path) or path.name.startswith("~$"):
                issues.append(_issue(path, "temporary_or_link_skipped"))
                continue
            extension = path.suffix.lower()
            if extension not in SUPPORTED_EXTENSIONS:
                code = (
                    "legacy_office_format_unsupported"
                    if extension in LEGACY_EXTENSIONS
                    else "ocr_not_supported"
                    if extension in IMAGE_EXTENSIONS
                    else "format_unsupported"
                )
                issues.append(_issue(path, code))
                continue
            key = os.path.normcase(str(path.resolve()))
            if key not in seen:
                seen.add(key)
                files.append(path.resolve())
    return sorted(files, key=lambda item: os.path.normcase(str(item))), issues


def _is_reparse_or_link(path: Path) -> bool:
    try:
        if path.is_symlink():
            return True
        is_junction = getattr(path, "is_junction", None)
        return bool(callable(is_junction) and is_junction())
    except OSError:
        return True


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _extract(path: Path) -> dict[str, Any]:
    extension = path.suffix.lower()
    try:
        if extension == ".pdf":
            return _extract_pdf(path)
        if extension == ".docx":
            return _extract_docx(path)
        if extension == ".xlsx":
            return _extract_xlsx(path)
        if extension == ".pptx":
            return _extract_pptx(path)
        return _extract_text(path)
    except Exception as exc:  # parser errors are evidence gaps, never process crashes
        name = type(exc).__name__.lower()
        code = "document_encrypted" if "password" in name or "encrypted" in str(exc).lower() else "document_parse_failed"
        return {"status": "unparseable", "issue": code, "chunks": []}


def _extract_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        return {"status": "unparseable", "issue": "document_encrypted", "chunks": []}
    chunks = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            chunks.append(_chunk(text, {"kind": "pdf_page", "page": page_number}))
    if not chunks:
        return {"status": "unparseable", "issue": "pdf_text_layer_unavailable", "chunks": []}
    return {"status": "parsed", "chunks": chunks}


def _extract_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    document = Document(str(path))
    chunks: list[dict[str, Any]] = []
    heading = ""
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        if str(paragraph.style.name).lower().startswith("heading"):
            heading = text
        chunks.append(
            _chunk(
                text,
                {"kind": "docx_paragraph", "paragraph": index, "heading": heading},
            )
        )
    for table_index, table in enumerate(document.tables, start=1):
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        text = "\n".join(row for row in rows if row.strip()).strip()
        if text:
            chunks.append(_chunk(text, {"kind": "docx_table", "table": table_index}))
    return _parsed_or_empty(chunks)


def _extract_xlsx(path: Path) -> dict[str, Any]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True, keep_links=False)
    chunks: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        batch: list[str] = []
        start_row = 1
        end_row = 0
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            values = ["" if value is None else str(value) for value in row]
            if not any(value.strip() for value in values):
                continue
            if not batch:
                start_row = row_index
            batch.append("\t".join(values))
            end_row = row_index
            if len(batch) >= 100:
                chunks.append(
                    _chunk(
                        "\n".join(batch),
                        {"kind": "xlsx_range", "sheet": sheet.title, "rows": [start_row, end_row]},
                    )
                )
                batch = []
        if batch:
            chunks.append(
                _chunk(
                    "\n".join(batch),
                    {"kind": "xlsx_range", "sheet": sheet.title, "rows": [start_row, end_row]},
                )
            )
    workbook.close()
    return _parsed_or_empty(chunks)


def _extract_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    chunks = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        values: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    values.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    values.append("\t".join(cell.text.strip() for cell in row.cells))
        text = "\n".join(values).strip()
        if text:
            chunks.append(_chunk(text, {"kind": "pptx_slide", "slide": slide_number}))
    return _parsed_or_empty(chunks)


def _extract_text(path: Path) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8-sig")
    if path.suffix.lower() == ".csv":
        rows = list(csv.reader(text.splitlines()))
        lines = ["\t".join(row) for row in rows]
    else:
        lines = text.splitlines()
    chunks = []
    for start in range(0, len(lines), 200):
        segment = "\n".join(lines[start : start + 200]).strip()
        if segment:
            chunks.append(
                _chunk(
                    segment,
                    {"kind": "text_lines", "lines": [start + 1, min(len(lines), start + 200)]},
                )
            )
    return _parsed_or_empty(chunks)


def _parsed_or_empty(chunks: list[dict[str, Any]]) -> dict[str, Any]:
    return (
        {"status": "parsed", "chunks": chunks}
        if chunks
        else {"status": "unparseable", "issue": "document_empty", "chunks": []}
    )


def _chunk(text: str, locator: dict[str, Any]) -> dict[str, Any]:
    content = text.strip()
    return {
        "chunk_id": "chunk_" + hashlib.sha256(
            (json.dumps(locator, sort_keys=True) + "\n" + content).encode("utf-8")
        ).hexdigest()[:20],
        "locator": locator,
        "text": content,
        "sha256": hashlib.sha256(content.encode("utf-8")).hexdigest(),
    }


def _issue(path: Path, code: str, **detail: Any) -> dict[str, Any]:
    return {"name": path.name, "path": str(path), "code": code, **detail}
